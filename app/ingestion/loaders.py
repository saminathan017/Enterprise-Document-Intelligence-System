"""
Document loaders — complete format support for 2026 enterprise stack.

Supported formats:
  PDF   — pdfplumber (structured table extraction) + PyPDF2 fallback
  PDF*  — GPT-4o Vision OCR for scanned / image-heavy pages
  TXT   — UTF-8 / latin-1
  MD    — Markdown
  DOCX  — python-docx
  XLSX/XLS/CSV — openpyxl / pandas → markdown tables
  PPTX  — python-pptx (slide text + speaker notes)
"""

from __future__ import annotations

import io
import base64
from abc import ABC, abstractmethod
from pathlib import Path
from typing import Dict, Any, List, Optional
from app.config import settings

# ── stdlib fallbacks so server starts even if optional deps are missing ────────
try:
    import pdfplumber
    _PDFPLUMBER = True
except ImportError:
    _PDFPLUMBER = False

try:
    import PyPDF2
    _PYPDF2 = True
except ImportError:
    _PYPDF2 = False

try:
    import docx as _docx
    _DOCX = True
except ImportError:
    _DOCX = False

try:
    import openpyxl as _openpyxl  # noqa: F401 — pandas xlsx backend
    import pandas  # noqa: F401 — presence check only; re-imported inside method
    _EXCEL = True
except ImportError:
    _EXCEL = False

try:
    from pptx import Presentation as _Presentation
    _PPTX = True
except ImportError:
    _PPTX = False


# ── helpers ───────────────────────────────────────────────────────────────────

def _table_to_markdown(table: List[List[Optional[str]]]) -> str:
    """Convert a pdfplumber / openpyxl table to GitHub Markdown."""
    if not table:
        return ""
    rows = [[str(c or "") for c in row] for row in table]
    if not rows:
        return ""
    header = "| " + " | ".join(rows[0]) + " |"
    sep    = "| " + " | ".join("---" for _ in rows[0]) + " |"
    body   = "\n".join("| " + " | ".join(r) + " |" for r in rows[1:])
    return f"{header}\n{sep}\n{body}"


def _decode_text(content: bytes) -> str:
    for enc in ("utf-8", "utf-16", "latin-1"):
        try:
            return content.decode(enc)
        except (UnicodeDecodeError, LookupError):
            continue
    return content.decode("latin-1", errors="replace")


# ── base class ─────────────────────────────────────────────────────────────────

class DocumentLoader(ABC):
    @abstractmethod
    def load(
        self,
        content: bytes,
        filename: str,
        metadata: Dict[str, Any] = None,
    ) -> Dict[str, Any]:
        ...


# ── PDF loader (pdfplumber + table extraction + Vision OCR fallback) ──────────

class PDFLoader(DocumentLoader):
    """
    PDF extraction strategy (in order):
      1. pdfplumber  — structured text + tables as Markdown
      2. PyPDF2      — plain text fallback
      3. GPT-4o Vision — scanned / zero-text pages
    """

    MIN_TEXT_CHARS_PER_PAGE = 50  # below this → treat as scanned

    def load(
        self,
        content: bytes,
        filename: str,
        metadata: Dict[str, Any] = None,
    ) -> Dict[str, Any]:
        page_texts: List[str] = []
        page_count = 0
        scanned_pages: List[int] = []

        if _PDFPLUMBER:
            try:
                with pdfplumber.open(io.BytesIO(content)) as pdf:
                    page_count = len(pdf.pages)
                    for i, page in enumerate(pdf.pages):
                        parts: List[str] = []

                        # Extract tables first (preserve structure)
                        tables = page.extract_tables() or []
                        table_bboxes = [t.bbox for t in (page.find_tables() or [])]

                        for tbl in tables:
                            md = _table_to_markdown(tbl)
                            if md:
                                parts.append(f"\n\n[TABLE — page {i+1}]\n{md}\n")

                        # Extract remaining text excluding table regions
                        if table_bboxes:
                            non_table = page
                            for bbox in table_bboxes:
                                try:
                                    non_table = non_table.outside_bbox(bbox)
                                except Exception:
                                    pass
                            raw_text = non_table.extract_text() or ""
                        else:
                            raw_text = page.extract_text() or ""

                        if raw_text.strip():
                            parts.append(raw_text)

                        combined = "\n".join(parts).strip()
                        if len(combined) < self.MIN_TEXT_CHARS_PER_PAGE:
                            scanned_pages.append(i)
                        else:
                            page_texts.append(combined)
            except Exception:
                pass  # fall through to PyPDF2

        if not page_texts and _PYPDF2:
            try:
                reader = PyPDF2.PdfReader(io.BytesIO(content))
                page_count = page_count or len(reader.pages)
                for page in reader.pages:
                    t = page.extract_text() or ""
                    page_texts.append(t)
            except Exception:
                pass

        # GPT-4o Vision for scanned pages
        if scanned_pages:
            vision_texts = self._ocr_pages(content, scanned_pages)
            page_texts.extend(vision_texts)

        full_text = "\n\n".join(t for t in page_texts if t.strip())

        base_meta: Dict[str, Any] = {
            "source":     filename,
            "file_type":  "pdf",
            "page_count": page_count,
            "has_tables": any("[TABLE" in t for t in page_texts),
            "scanned_pages": len(scanned_pages),
            **(metadata or {}),
        }

        # Try to extract PDF metadata
        if _PYPDF2:
            try:
                reader = PyPDF2.PdfReader(io.BytesIO(content))
                if reader.metadata:
                    if reader.metadata.title:
                        base_meta["title"] = reader.metadata.title
                    if reader.metadata.author:
                        base_meta["author"] = reader.metadata.author
            except Exception:
                pass

        return {"text": full_text, "metadata": base_meta}

    # ── Vision OCR ────────────────────────────────────────────────────────────

    @staticmethod
    def _ocr_pages(pdf_bytes: bytes, page_indices: List[int]) -> List[str]:
        """
        Use GPT-4o Vision to OCR scanned pages.
        Converts each page to PNG via pdf2image (optional dep) then calls API.
        Falls back silently if dependencies are missing.
        """
        try:
            from pdf2image import convert_from_bytes
        except ImportError:
            return []

        try:
            from openai import OpenAI
            from app.config import settings
            client = OpenAI(api_key=settings.openai_api_key)
        except Exception:
            return []

        results: List[str] = []
        try:
            images = convert_from_bytes(pdf_bytes, dpi=150)
            for i in page_indices:
                if i >= len(images):
                    continue
                buf = io.BytesIO()
                images[i].save(buf, format="PNG")
                b64 = base64.b64encode(buf.getvalue()).decode()
                resp = client.chat.completions.create(
                    model=settings.ocr_model,
                    messages=[{
                        "role": "user",
                        "content": [
                            {
                                "type": "text",
                                "text": (
                                    "Extract ALL text from this scanned document page. "
                                    "Preserve tables as Markdown. Return only the extracted text."
                                ),
                            },
                            {
                                "type": "image_url",
                                "image_url": {"url": f"data:image/png;base64,{b64}", "detail": "high"},
                            },
                        ],
                    }],
                    max_tokens=2000,
                )
                results.append(resp.choices[0].message.content or "")
        except Exception:
            pass

        return results


# ── Plain text ─────────────────────────────────────────────────────────────────

class TextLoader(DocumentLoader):
    def load(self, content: bytes, filename: str, metadata: Dict[str, Any] = None) -> Dict[str, Any]:
        text = _decode_text(content)
        return {
            "text": text,
            "metadata": {
                "source": filename,
                "file_type": "txt",
                "char_count": len(text),
                **(metadata or {}),
            },
        }


# ── Markdown ───────────────────────────────────────────────────────────────────

class MarkdownLoader(DocumentLoader):
    def load(self, content: bytes, filename: str, metadata: Dict[str, Any] = None) -> Dict[str, Any]:
        text = _decode_text(content)
        return {
            "text": text,
            "metadata": {
                "source": filename,
                "file_type": "markdown",
                "char_count": len(text),
                **(metadata or {}),
            },
        }


# ── DOCX ───────────────────────────────────────────────────────────────────────

class DocxLoader(DocumentLoader):
    def load(self, content: bytes, filename: str, metadata: Dict[str, Any] = None) -> Dict[str, Any]:
        if not _DOCX:
            raise ImportError("python-docx not installed. Run: pip install python-docx")

        doc = _docx.Document(io.BytesIO(content))
        parts: List[str] = []

        for elem in doc.element.body:
            tag = elem.tag.split("}")[-1] if "}" in elem.tag else elem.tag

            if tag == "p":
                # Regular paragraph
                text = "".join(r.text for r in elem.findall(".//{http://schemas.openxmlformats.org/wordprocessingml/2006/main}t"))
                if text.strip():
                    parts.append(text)

            elif tag == "tbl":
                # Table → Markdown
                rows: List[List[str]] = []
                ns = "http://schemas.openxmlformats.org/wordprocessingml/2006/main"
                for tr in elem.findall(f"{{{ns}}}tr"):
                    row: List[str] = []
                    for tc in tr.findall(f"{{{ns}}}tc"):
                        cell_text = "".join(
                            t.text or "" for t in tc.findall(f".//{{{ns}}}t")
                        )
                        row.append(cell_text.strip())
                    rows.append(row)
                md = _table_to_markdown(rows)
                if md:
                    parts.append(f"\n[TABLE]\n{md}\n")

        full_text = "\n\n".join(parts)
        return {
            "text": full_text,
            "metadata": {
                "source": filename,
                "file_type": "docx",
                "char_count": len(full_text),
                **(metadata or {}),
            },
        }


# ── Excel / CSV ────────────────────────────────────────────────────────────────

class ExcelLoader(DocumentLoader):
    def load(self, content: bytes, filename: str, metadata: Dict[str, Any] = None) -> Dict[str, Any]:
        ext = Path(filename).suffix.lower()

        if ext == ".csv":
            text = _decode_text(content)
            # Convert CSV to Markdown table
            import csv as _csv
            reader = list(_csv.reader(io.StringIO(text)))
            md = _table_to_markdown(reader) if reader else text
            full_text = f"[CSV: {filename}]\n{md}"
        else:
            if not _EXCEL:
                raise ImportError("openpyxl not installed. Run: pip install openpyxl pandas")
            import pandas as pd
            xl = pd.ExcelFile(io.BytesIO(content))
            parts: List[str] = []
            for sheet in xl.sheet_names:
                df = xl.parse(sheet)
                df = df.fillna("")
                rows = [df.columns.tolist()] + df.values.tolist()
                md = _table_to_markdown([[str(c) for c in r] for r in rows])
                parts.append(f"[SHEET: {sheet}]\n{md}")
            full_text = "\n\n".join(parts)

        return {
            "text": full_text,
            "metadata": {
                "source": filename,
                "file_type": ext.lstrip("."),
                "has_tables": True,
                **(metadata or {}),
            },
        }


# ── PowerPoint ─────────────────────────────────────────────────────────────────

class PptxLoader(DocumentLoader):
    def load(self, content: bytes, filename: str, metadata: Dict[str, Any] = None) -> Dict[str, Any]:
        if not _PPTX:
            raise ImportError("python-pptx not installed. Run: pip install python-pptx")

        prs = _Presentation(io.BytesIO(content))
        parts: List[str] = []

        for i, slide in enumerate(prs.slides, 1):
            slide_parts: List[str] = [f"[SLIDE {i}]"]
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        t = "".join(run.text for run in para.runs).strip()
                        if t:
                            slide_parts.append(t)
                if shape.has_table:
                    rows = [
                        [cell.text.strip() for cell in row.cells]
                        for row in shape.table.rows
                    ]
                    md = _table_to_markdown(rows)
                    if md:
                        slide_parts.append(f"\n[TABLE]\n{md}\n")

            # Speaker notes
            if slide.has_notes_slide:
                notes_tf = slide.notes_slide.notes_text_frame
                if notes_tf:
                    notes = notes_tf.text.strip()
                    if notes:
                        slide_parts.append(f"[NOTES] {notes}")

            parts.append("\n".join(slide_parts))

        full_text = "\n\n".join(parts)
        return {
            "text": full_text,
            "metadata": {
                "source": filename,
                "file_type": "pptx",
                "slide_count": len(prs.slides),
                **(metadata or {}),
            },
        }


# ── Factory ────────────────────────────────────────────────────────────────────

class DocumentLoaderFactory:
    _loaders = {
        ".pdf":      PDFLoader,
        ".txt":      TextLoader,
        ".md":       MarkdownLoader,
        ".markdown": MarkdownLoader,
        ".docx":     DocxLoader,
        ".doc":      DocxLoader,
        ".xlsx":     ExcelLoader,
        ".xls":      ExcelLoader,
        ".csv":      ExcelLoader,
        ".pptx":     PptxLoader,
        ".ppt":      PptxLoader,
    }

    @classmethod
    def get_loader(cls, filename: str) -> DocumentLoader:
        ext = Path(filename).suffix.lower()
        loader_cls = cls._loaders.get(ext)
        if not loader_cls:
            supported = ", ".join(sorted(cls._loaders.keys()))
            raise ValueError(f"Unsupported file type: {ext}. Supported: {supported}")
        return loader_cls()

    @classmethod
    def supported_extensions(cls) -> List[str]:
        return sorted(cls._loaders.keys())
